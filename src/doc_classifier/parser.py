from __future__ import annotations

import logging
import re
import shutil
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

SUPPORTED_TEXT = {".pdf", ".docx", ".doc", ".txt"}
SUPPORTED_IMAGE = {".jpg", ".jpeg", ".jpe", ".png", ".bmp", ".tiff", ".gif"}
SUPPORTED_OFFICE = {
    ".xlsx", ".xls", ".csv", ".tsv", ".ods", ".odt",
    ".pptx", ".ppt", ".pot", ".xlb",
}
MEDIA_EXTENSIONS = {".mp3", ".wav", ".wma", ".ogg", ".au", ".mp2", ".mid", ".ram", ".rm"}
MAX_CHARS = 2000
GIBBERISH_THRESHOLD = 0.3

_TESSERACT_PATHS = [
    Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe"),
    Path(r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"),
]

URL_PATTERN = re.compile(r"https?://[^\s]+", re.IGNORECASE)
FONT_NAMES = {"arial", "calibri", "times new roman", "verdana", "tahoma", "century gothic",
              "garamond", "gill sans", "zapfhumnst", "courier", "helvetica", "symbol"}


def _strip_urls(text: str) -> str:
    return URL_PATTERN.sub("", text)


def _is_boilerplate(text: str) -> bool:
    lower = text.lower()
    if "online service" in lower and "e-mail" in lower:
        return True
    if "please e-mail" in lower and "http" in lower:
        return True
    return False


def _configure_tesseract() -> bool:
    try:
        import pytesseract
    except ImportError:
        logger.warning("pytesseract not installed")
        return False

    current = pytesseract.pytesseract.tesseract_cmd
    if current not in ("tesseract", ""):
        return True

    found = shutil.which("tesseract")
    if not found:
        found = next((str(p) for p in _TESSERACT_PATHS if p.exists()), None)
    if found:
        pytesseract.pytesseract.tesseract_cmd = found
        logger.info("Using tesseract at %s", found)
        return True

    logger.warning(
        "tesseract executable not found in PATH or %s", ", ".join(str(p) for p in _TESSERACT_PATHS)
    )
    return False


def _extract_pdf_metadata(file_path: Path) -> str | None:
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            meta = pdf.metadata
            if meta:
                for key in ("title", "Subject", "subject"):
                    val = meta.get(key)
                    if val and isinstance(val, str) and val.strip():
                        cleaned = val.strip()
                        if len(cleaned) > 5 and len(cleaned) < 200:
                            return cleaned
    except Exception:
        pass
    return None


def parse_pdf(file_path: Path) -> str:
    import pdfplumber

    text_parts: list[str] = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages[:5]:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
            if sum(len(t) for t in text_parts) >= MAX_CHARS:
                break

    combined = "\n".join(text_parts)
    if not combined.strip():
        return _ocr_pdf_fallback(file_path)
    return combined[:MAX_CHARS]


def _render_pdf_pages(file_path: Path, max_pages: int = 3):
    try:
        import pdf2image
        return pdf2image.convert_from_path(
            str(file_path), first_page=1, last_page=max_pages
        )
    except Exception as exc:
        logger.info("pdf2image unavailable (%s) - using pdfplumber renderer", exc)

    import pdfplumber
    with pdfplumber.open(file_path) as pdf:
        return [page.to_image().original for page in pdf.pages[:max_pages]]


def _ocr_pdf_fallback(file_path: Path) -> str:
    if not _configure_tesseract():
        logger.warning("Tesseract not configured - skipping OCR fallback for %s", file_path)
        return ""
    try:
        import pytesseract
        images = _render_pdf_pages(file_path)
        text_parts = []
        conf_sum = 0.0
        conf_count = 0
        for img in images:
            data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
            for i, conf_str in enumerate(data["conf"]):
                try:
                    conf = int(conf_str)
                    if conf > 0:
                        conf_sum += conf
                        conf_count += 1
                        word = (data["text"][i] or "").strip()
                        if word:
                            text_parts.append(word)
                except ValueError:
                    pass
        avg_conf = conf_sum / conf_count if conf_count else 0.0
        if avg_conf < 20.0:
            logger.info(
                "OCR confidence %.0f%% too low for %s, fallback to filename",
                avg_conf, file_path.name,
            )
            return ""
        return "\n".join(text_parts)[:MAX_CHARS]
    except Exception as exc:
        logger.warning("OCR fallback failed for %s: %s", file_path, exc)
        return ""


def parse_docx(file_path: Path) -> str:
    import docx

    doc = docx.Document(str(file_path))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return text[:MAX_CHARS]


def parse_txt(file_path: Path) -> str:
    encodings = ["utf-8", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            return file_path.read_text(encoding=enc)[:MAX_CHARS]
        except (UnicodeDecodeError, LookupError):
            continue
    return ""


def parse_csv(file_path: Path) -> str:
    import csv

    text_parts: list[str] = []
    try:
        with open(file_path, "r", encoding="utf-8-sig", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                line = ", ".join(c.strip() for c in row if c.strip())
                if line:
                    text_parts.append(line)
                if sum(len(t) for t in text_parts) >= MAX_CHARS:
                    break
    except Exception as exc:
        logger.error("Error parsing %s: %s", file_path, exc)
        return ""
    return "\n".join(text_parts)[:MAX_CHARS]


def parse_tsv(file_path: Path) -> str:
    import csv

    text_parts: list[str] = []
    try:
        with open(file_path, "r", encoding="utf-8-sig", errors="replace") as f:
            reader = csv.reader(f, delimiter="\t")
            for row in reader:
                line = ", ".join(c.strip() for c in row if c.strip())
                if line:
                    text_parts.append(line)
                if sum(len(t) for t in text_parts) >= MAX_CHARS:
                    break
    except Exception as exc:
        logger.error("Error parsing %s: %s", file_path, exc)
        return ""
    return "\n".join(text_parts)[:MAX_CHARS]


def parse_xlsx(file_path: Path) -> str:
    import openpyxl

    text_parts: list[str] = []
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames[:3]:
            ws = wb[sheet_name]
            text_parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                line = ", ".join(cells)
                if line:
                    text_parts.append(line)
                if sum(len(t) for t in text_parts) >= MAX_CHARS:
                    break
            if sum(len(t) for t in text_parts) >= MAX_CHARS:
                break
        wb.close()
    except Exception as exc:
        logger.error("Error parsing %s: %s", file_path, exc)
        return ""
    return "\n".join(text_parts)[:MAX_CHARS]


def _parse_xls_as_html(file_path: Path) -> Optional[str]:
    try:
        from bs4 import BeautifulSoup
        content = file_path.read_text(encoding="utf-8-sig", errors="replace")
        soup = BeautifulSoup(content, "html.parser")
        tables = soup.find_all("table")
        if not tables:
            return None
        text_parts: list[str] = []
        for table in tables[:3]:
            rows = table.find_all("tr")
            for row in rows[:100]:
                cells = row.find_all(["td", "th"])
                cell_texts = [c.get_text(strip=True) for c in cells if c.get_text(strip=True)]
                line = ", ".join(cell_texts)
                if line:
                    text_parts.append(line)
                if sum(len(t) for t in text_parts) >= MAX_CHARS:
                    break
            if sum(len(t) for t in text_parts) >= MAX_CHARS:
                break
        return "\n".join(text_parts)[:MAX_CHARS]
    except Exception as exc:
        logger.warning("HTML fallback for .xls failed: %s", exc)
        return None


def parse_xls(file_path: Path) -> str:
    import xlrd

    try:
        wb = xlrd.open_workbook(str(file_path), on_demand=True)
    except Exception as exc:
        logger.info("xlrd failed for %s: %s - trying HTML table fallback", file_path, exc)
        html_text = _parse_xls_as_html(file_path)
        if html_text:
            return html_text
        logger.warning("HTML fallback also failed for %s", file_path)
        return ""

    text_parts: list[str] = []
    try:
        for sheet_idx in range(min(wb.nsheets, 3)):
            sheet = wb.sheet_by_index(sheet_idx)
            text_parts.append(f"[Sheet: {sheet.name}]")
            for row_idx in range(min(sheet.nrows, 100)):
                cells = [
                    str(sheet.cell_value(row_idx, col_idx)).strip()
                    for col_idx in range(sheet.ncols)
                ]
                line = ", ".join(c for c in cells if c)
                if line:
                    text_parts.append(line)
                if sum(len(t) for t in text_parts) >= MAX_CHARS:
                    break
            if sum(len(t) for t in text_parts) >= MAX_CHARS:
                break
    finally:
        wb.release_resources()
    return "\n".join(text_parts)[:MAX_CHARS]


def parse_ods(file_path: Path) -> str:
    import odf.opendocument
    import odf.text as odf_text

    doc = odf.opendocument.load(file_path)
    text_parts: list[str] = []
    for elem in doc.getElementsByType(odf_text.P):
        if elem.childNodes:
            text = "".join(node.data for node in elem.childNodes if hasattr(node, "data"))
            if text.strip():
                text_parts.append(text.strip())
            if sum(len(t) for t in text_parts) >= MAX_CHARS:
                break
    return "\n".join(text_parts)[:MAX_CHARS]


def parse_odt(file_path: Path) -> str:
    import odf.opendocument
    import odf.text as odf_text

    doc = odf.opendocument.load(file_path)
    text_parts: list[str] = []
    for elem in doc.getElementsByType(odf_text.P):
        if elem.childNodes:
            text = "".join(node.data for node in elem.childNodes if hasattr(node, "data"))
            if text.strip():
                text_parts.append(text.strip())
            if sum(len(t) for t in text_parts) >= MAX_CHARS:
                break
    return "\n".join(text_parts)[:MAX_CHARS]


def _parse_ppt_binary(file_path: Path) -> Optional[str]:
    try:
        import olefile

        ole = olefile.OleFileIO(str(file_path))
        text_parts: list[str] = []
        try:
            for stream_name in ole.listdir():
                name = "/".join(stream_name)
                if "PowerPoint Document" in name:
                    data = ole.openstream(stream_name).read()
                    text = data.decode("utf-16-le", errors="replace")
                    cleaned = []
                    for c in text:
                        if c.isprintable() or c in "\n\r\t":
                            if c.isascii() or c in " \n\r\t":
                                cleaned.append(c)
                            elif c.isalpha() and ord(c) < 256:
                                cleaned.append(c)
                            else:
                                cleaned.append(" ")
                        elif c in (" ", "\x00"):
                            cleaned.append(" ")
                    clean_str = "".join(cleaned)
                    clean_str = _strip_urls(clean_str)
                    words = [w.strip() for w in clean_str.split() if len(w.strip()) > 2]
                    words = [w for w in words if w.lower() not in FONT_NAMES]
                    content = " ".join(words)
                    if content:
                        text_parts.append(content[:MAX_CHARS])
        finally:
            ole.close()
        result = "\n".join(text_parts)[:MAX_CHARS] if text_parts else None
        if result and len(result) > 20:
            alpha = sum(1 for c in result if c.isalpha())
            if alpha / len(result) < 0.3:
                return None
        return result
    except ImportError:
        logger.warning("olefile not installed - cannot parse binary .ppt")
        return None
    except Exception as exc:
        logger.warning("Binary .ppt parsing failed for %s: %s", file_path, exc)
        return None


def parse_pptx(file_path: Path) -> str:
    ext = file_path.suffix.lower()
    if ext == ".ppt":
        binary_text = _parse_ppt_binary(file_path)
        if binary_text:
            return binary_text
        return ""

    from pptx import Presentation

    prs = Presentation(str(file_path))
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
            if sum(len(t) for t in text_parts) >= MAX_CHARS:
                break
        if sum(len(t) for t in text_parts) >= MAX_CHARS:
            break
    return "\n".join(text_parts)[:MAX_CHARS]


def _is_high_quality_text(text: str) -> bool:
    if not text or len(text) < 20:
        return False
    alpha_chars = sum(1 for c in text if c.isalpha())
    total_chars = len(text.strip())
    if total_chars == 0:
        return False
    ratio = alpha_chars / total_chars
    return ratio >= GIBBERISH_THRESHOLD


def parse_office(file_path: Path) -> str:
    ext = file_path.suffix.lower()
    if ext == ".csv":
        return parse_csv(file_path)
    if ext == ".tsv":
        return parse_tsv(file_path)
    if ext == ".xlsx":
        return parse_xlsx(file_path)
    if ext == ".xls":
        return parse_xls(file_path)
    if ext == ".ods":
        return parse_ods(file_path)
    if ext == ".odt":
        return parse_odt(file_path)
    if ext in (".ppt", ".pptx", ".pot"):
        return parse_pptx(file_path)
    if ext == ".xlb":
        return ""
    raise ValueError(f"Unsupported office extension: {ext}")


def _ocr_image_data(file_path: Path) -> tuple[str, float]:
    """OCR image and return (text, avg_confidence)."""
    import pytesseract
    from PIL import Image

    img = Image.open(file_path)
    data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
    confidences = []
    text_parts = []
    for i, conf_str in enumerate(data["conf"]):
        try:
            conf = int(conf_str)
            if conf > 0:
                confidences.append(conf)
                word = (data["text"][i] or "").strip()
                if word:
                    text_parts.append(word)
        except ValueError:
            pass

    text = " ".join(text_parts)
    avg_conf = sum(confidences) / len(confidences) if confidences else 0.0
    return text, avg_conf


def parse_image(file_path: Path) -> str:
    if not _configure_tesseract():
        logger.warning("Tesseract not configured - cannot OCR %s", file_path)
        return ""
    try:
        text, confidence = _ocr_image_data(file_path)
        if confidence < 20.0:
            logger.info(
                "OCR confidence %.0f%% too low for %s, fallback to filename",
                confidence, file_path.name,
            )
            return ""
        return text.strip()[:MAX_CHARS]
    except ImportError:
        logger.warning("pytesseract/Pillow not installed - cannot OCR %s", file_path)
        return ""
    except Exception as exc:
        logger.warning("OCR failed for %s: %s", file_path, exc)
        return ""


def extract_text_metadata(file_path: str | Path) -> tuple[str | None, str | None]:
    path = Path(file_path)
    ext = path.suffix.lower()

    title_meta = None
    text = None

    if ext == ".pdf":
        title_meta = _extract_pdf_metadata(path)

    try:
        text = extract_text(path)
    except Exception:
        pass

    if text:
        text = _strip_urls(text)
        lines = [
            line.strip() for line in text.split("\n")
            if line.strip() and len(line.strip()) > 15
        ]
        lines = [line for line in lines if not _is_boilerplate(line)]
        text = "\n".join(lines)

    return title_meta, text


def extract_text(file_path: str | Path) -> Optional[str]:
    path = Path(file_path)
    if not path.exists():
        logger.error("File not found: %s", path)
        return None

    ext = path.suffix.lower()
    text = ""

    try:
        if ext == ".pdf":
            text = parse_pdf(path)
        elif ext in (".docx", ".doc"):
            text = parse_docx(path)
        elif ext == ".txt":
            text = parse_txt(path)
        elif ext in SUPPORTED_OFFICE:
            text = parse_office(path)
        elif ext in SUPPORTED_IMAGE:
            text = parse_image(path)
        else:
            logger.warning("Unsupported file type: %s", ext)
            return None
    except Exception as exc:
        logger.error("Error parsing %s: %s", path, exc)
        return None

    if not text or not text.strip():
        return None

    text = text.strip()

    if not _is_high_quality_text(text):
        return None

    return text


def get_supported_extensions() -> set[str]:
    return SUPPORTED_TEXT | SUPPORTED_IMAGE | SUPPORTED_OFFICE


def is_media_file(file_path: str | Path) -> bool:
    ext = Path(file_path).suffix.lower()
    return ext in MEDIA_EXTENSIONS
